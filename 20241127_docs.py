#%%
#pdf
# ms docs, hwp, 메모장, 서식(text, 표, 글씨)
# mistral ai LLM -> langchain (openAI, mistralAI) API KEY를 사용해야 각각의 LLM에 접근할 수 있음.

#%%

from langchain_core.prompts import ChatPromptTemplate
from langchain_mistralai import ChatMistralAI
import os



api_key = "sM5AM2CWw2X1C33oQQSjHy4WsTjSZm81"


chat = ChatMistralAI(
    model = "mistral-large-latest",
    temperature = 0,
    max_retries = 2,
    api_key=api_key, 
)    

message = [("system","You are a helpful assistant that translates English to Korean. Translate the user sentence."), 
           ("human", "What a nice weather!!")
           ]

ai_msg = chat.invoke(message)

print(ai_msg.content)

#%%
# 초기 메시지 설정
message = [("system", "You are a helpful assistant that translates English to Korean. Translate the user sentence.")]

# 사용자 입력을 받아 대화를 이어가는 루프
while True:
    # 사용자 입력 받기
    user_input = input("You: ")

    # 종료 조건 설정
    if user_input.lower() in ["exit", "quit", "bye"]:
        print("Goodbye!")
        break

    # 사용자 입력을 메시지 리스트에 추가
    message.append(("human", user_input))

    # 모델에 메시지 전달하여 응답 받기
    ai_msg = chat.invoke(message)

    # 모델의 응답을 출력
    print(f"AI: {ai_msg.content}")

    # 모델의 응답을 메시지 리스트에 추가
    message.append(("ai", ai_msg.content))

    
#%%
from langchain_core.prompts import ChatPromptTemplate
from langchain_mistralai import ChatMistralAI
import os



api_key = "sM5AM2CWw2X1C33oQQSjHy4WsTjSZm81"


chat = ChatMistralAI(
    model = "mistral-large-latest",
    temperature = 0,
    max_retries = 2,
    api_key=api_key, 
)  

# 초기 메시지 설정
message = [("system", "You are a helpful assistant that translates English to Korean. Translate the user sentence.")]

# 사용자 입력을 받아 대화를 이어가는 루프
while True:
    # 사용자 입력 받기
    user_input = input("You: ")

    # 종료 조건 설정
    if user_input.lower() in ["exit", "quit", "bye"]:
        print("Goodbye!")
        break

    # 사용자 입력을 메시지 리스트에 추가
    message.append(("human", user_input))

    # 모델에 메시지 전달하여 응답 받기
    ai_msg = chat.invoke(message)

    # 모델의 응답을 출력
    print(f"AI: {ai_msg.content}")

    # 모델의 응답을 메시지 리스트에 추가
    message.append(("ai", ai_msg.content))
















# %%
# pip install langchain_community
# pip install pypdf
from langchain.document_loaders import PyPDFLoader





# %%
from langchain.document_loaders import PyPDFLoader
file_path ="./[이슈리포트 2022-2호] 혁신성장 정책금융 동향.pdf"
loader = PyPDFLoader(file_path)
#%%
contents = loader.load_and_split()
print("*"*10, contents)
print(contents[2].page_content)
# %%
# 문서에 있는 이미지 추출
# pip install rapidocr_onnxruntime
loader_img = PyPDFLoader(file_path, extract_images=True)

contents_img = loader_img.load_and_split()
# %%
# pip install pypdfium2
from langchain.document_loaders import PyPDFium2Loader
loader_total = PyPDFium2Loader(file_path)
contents_total = loader_total.load()
print(contents_total[3].page_content)
# %%

from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import CSVLoader
from langchain_community.document_loaders import PyPDFium2Loader
from langchain_community.document_loaders import UnstructuredPowerPointLoader

# CSV 파일 경로 지정
csv_path = './SP500_NOV2019_Hist.csv'

# CSVLoader 인스턴스 생성
loader = CSVLoader(file_path=csv_path)

# CSV 파일 로드
documents = loader.load()

print(documents)
# 로드된 데이터 출력
for doc in documents:
    print(doc.page_content)
    
    
documents[0]
documents[0].metadata
type(documents)   
# %%
import pandas as pd

# CSV 파일 경로 지정
csv_path = './SP500_NOV2019_Hist.csv'

# CSV 파일 읽기
df = pd.read_csv(csv_path)

# 데이터프레임 출력
print(df)

df
# %%
from docx import Document

def convert_txt_to_docx(txt_path, docx_path):
    # Read the text file
    with open(txt_path, 'r', encoding='utf-8') as file:
        text_content = file.read()

    # Create a new Document object
    doc = Document()

    # Add the text content to the Document
    doc.add_paragraph(text_content)

    # Save the Document as a .docx file
    doc.save(docx_path)

# Example usage
txt_path = './output/output.txt'
docx_path = './output/file.docx'

convert_txt_to_docx(txt_path, docx_path)

# %%
from langchain.document_loaders import Docx2txtLoader
from docx import Document

def convert_txt_to_docx(txt_path, docx_path):
    # Read the text file
    with open(txt_path, 'r', encoding='utf-8') as file:
        text_content = file.read()

    # Create a new Document object
    doc = Document()

    # Add the text content to the Document
    doc.add_paragraph(text_content)

    # Save the Document as a .docx file
    doc.save(docx_path)

def load_and_split(docx_path, txt_path, output_docx_path):
    # Load the .docx file using Docx2txtLoader
    loader = Docx2txtLoader(docx_path)
    documents = loader.load()

    # Combine text from all pages
    text_content = "\n".join([doc.page_content for doc in documents])

    # Save the text content to a .txt file
    with open(txt_path, 'w', encoding='utf-8') as file:
        file.write(text_content)

    # Convert the .txt file back to a .docx file
    convert_txt_to_docx(txt_path, output_docx_path)

# Example usage
docx_path = './output/file.docx'
txt_path = './output/text1.txt'
output_docx_path = './output/output_file.docx'

load_and_split(docx_path, txt_path, output_docx_path)

# %%
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

def convert_docx_to_pptx(docx_path, pptx_path):
    # Load the .docx file
    docx = DocxDocument(docx_path)

    # Create a new Presentation object
    prs = Presentation()

    # Iterate through the paragraphs in the .docx file
    for paragraph in docx.paragraphs:
        # Add a new slide for each paragraph
        slide_layout = prs.slide_layouts[1]  # Use the Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set the title and content of the slide
        title.text = "Slide Title"
        content.text = paragraph.text

    # Save the Presentation as a .pptx file
    prs.save(pptx_path)

# Example usage
docx_path = './output/file.docx'
pptx_path = './output/file.pptx'

convert_docx_to_pptx(docx_path, pptx_path)

# %%
from langchain_community.document_loaders import WebBaseLoader
import bs4


loaderw = WebBaseLoader('https://news.naver.com', encoding='UTF-8', 
                        bs_kwargs=dict(parse_only=bs4.SoupStrainer(class_ = ("comp_news_feed comp_news_none"))))
loaderw.requests_per_second = {'verify': False}
data = loaderw.load()

rslt = data[0].page_content
rslt.split("\n")
[item for item in data[0].page_content.split('\n') if item]
rslt = [item for item in data[0].page_content.split(
    '\n') if item and item not in ['구독', '영상']]
rslt




# %%
from langchain_community.document_loaders import WebBaseLoader
import bs4

def load_and_process_webpage(url):
    # Load the web page using WebBaseLoader
    loader = WebBaseLoader(url)
    loader.requests_per_second = {'verify': False}
    documents = loader.load()

    # Combine text from all pages
    text_content = "\n".join([doc.page_content for doc in documents])

    # Process the text content using BeautifulSoup
    soup = bs4.BeautifulSoup(text_content, 'html.parser')
    processed_text = soup.get_text()

    return processed_text

def save_text_to_file(text, file_path):
    # Save the processed text to a .txt file
    with open(file_path, 'w', encoding='utf-8') as file:
        file.write(text)

# Example usage
url = 'https://news.naver.com'
processed_text = load_and_process_webpage(url)
txt_file_path = './output/output.txt'
save_text_to_file(processed_text, txt_file_path)
print(f"Processed text has been saved to {txt_file_path}")

# %%
from langchain_core.prompts import ChatPromptTemplate
from langchain_mistralai import ChatMistralAI
import os
from langchain.schema import HumanMessage, AIMessage, SystemMessage
#from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain.callbacks import StreamingStdOutCallbackHandler


# API 키 가져오기
api_key = "sM5AM2CWw2X1C33oQQSjHy4WsTjSZm81"

# Mistral AI 모델 초기화
chat = ChatMistralAI(
    model="mistral-large-latest",
    temperature=0.5,
    max_retries=2,
    api_key=api_key,
    streaming=True,
    callbacks=[StreamingStdOutCallbackHandler()]
    
)

# 메시지 정의
message = [
    SystemMessage(
        "You are a helpful assistant that has expert knowledge about geography. Answer the user's question in {language}."
    ),
    AIMessage(
        "Hello, I'm a geography expert. My name is {name}!"
    ),
    HumanMessage(
        "How far is {ACity} from {BCity}? And also, how many hours does it take to fly to there?"
    )
]

# 프롬프트 템플릿 생성
template = ChatPromptTemplate.from_messages(
    [
        SystemMessage(
            "You are a helpful assistant that has expert knowledge about geography. Answer the user's question in {language}."
        ),
        AIMessage(
            "Hello, I'm a geography expert. My name is {name}!"
        ),
        HumanMessage(
            "How far is {ACity} from {BCity}? And also, how many hours does it take to fly there?"
        )
    ]
)

# 프롬프트 포맷팅
prompt = template.format(
    language="Korean",
    name="GeoExpert",
    ACity="LA",
    BCity="Seoul"
)

# 예측 요청
response = chat.predict(prompt)

# 예측된 메시지 출력
print(response)
# %%

chain = template | chat
chain.invoke(
    {
        "max_item":5,
        "question": "What are the top 5 countries with the highest GDP"
    }
)
# %%
chef_prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are a chef assistant. You will be asked to make a dish with easy to find ingredients. Do NOT reply with anything else. Also, always answer in Korean."),
        ("human", "I want to cook {menu}."),
    ]
)

chef_chain = chef_prompt  | chat

# %%
veg_chef_prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are a vegeterian chef specialist. You find alternative ingredients and explain their preparation. You don't modify the recipe. if there is no alternative for a good food just say you don't know how to replace it. and also, always answer in Korean."),
        ("human", "I want to know the vegeterian receipe.")
    ]
)

veg_chain = veg_chef_prompt | chat
# %%
#{"menu":"bibimbap"}
#{"recipe":"bulgogi"}




final_chain = {"recipe":chef_chain} | veg_chain
final_chain.invoke({"menu":"bibimbap"})

final_chain.invoke({"menu":"떡볶이", "receipe":"어묵"})

# %%
